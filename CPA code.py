from collections import abc
from pptx import Presentation
from pptx.util import Inches, Pt
#from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
import openpyxl


# Open the PowerPoint file
prs = Presentation()

#width1 = Inches(8.27)
#height1 = Inches(11.69)
wb = openpyxl.load_workbook("CPA.xlsx")
sheets = wb.sheetnames
sh1= wb['CPA']
sh2 = wb['Final CPA Content']
Name = sh1['A2'].value
Date = sh1['C2'].value
Signature_Here = "Signature Here"

# Change the size of the slides to 1920x1080 pixels 2480 3508
prs.slide_width = Inches(8.27)
prs.slide_height = Inches(11.69)
blank_slide_layout = prs.slide_layouts[6]

#Slide1
slide1 = prs.slides.add_slide(blank_slide_layout)
#slide1.width = Inches(8.27)
#slide1.height = Inches(11.69)

img_path1 = "Slide1img.jpg"
left = Inches(0)
top = Inches(0)

pic1 = slide1.shapes.add_picture(img_path1, left, top)

name_box = slide1.shapes.add_textbox(left=Inches(0.5530),top = Inches(9.59), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Name

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(10) 

date_box = slide1.shapes.add_textbox(left=Inches(0.5530),top = Inches(10.09), width=Inches(1.5), height=Inches(0.6))
date_frame = date_box.text_frame
date_frame.clear()
p = date_frame.paragraphs[0]
run = p.add_run()
run.text = Date

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(10)

#Slide2
slide2 = prs.slides.add_slide(blank_slide_layout)
#slide1.width = Inches(8.27)
#slide1.height = Inches(11.69)

img_path2 = "Slide2img.jpg"
left = Inches(0)
top = Inches(0)

pic2 = slide2.shapes.add_picture(img_path2, left, top)

#Slide3
slide3 = prs.slides.add_slide(blank_slide_layout)
#slide1.width = Inches(8.27)
#slide1.height = Inches(11.69)

img_path3 = "Slide3img.jpg"
left = Inches(0)
top = Inches(0)

pic3 = slide3.shapes.add_picture(img_path3, left, top)


#Slide4
slide4 = prs.slides.add_slide(blank_slide_layout)
img_path4 = "Slide4img.jpg"
left = Inches(0)
top = Inches(0)

pic4 = slide4.shapes.add_picture(img_path4, left, top)
#Slide 4 name
name_box = slide4.shapes.add_textbox(left=Inches(6.7),top = Inches(0.5), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Name

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(10) 


# Excel data fetching

dataa = sh1['D2'].value
datac = sh1['E2'].value
datae = sh1['F2'].value
datas = sh1['G2'].value
Advice_score = str(dataa)
Criticize_Score = str(datac)
Empathy_Score = str(datae)
Search_Score = str(datas)

#Slide 4 Advice score
name_box = slide4.shapes.add_textbox(left=Inches(1),top = Inches(2.75), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Advice_score

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(10) 
#Slide 4 Cri Score
name_box = slide4.shapes.add_textbox(left=Inches(3),top = Inches(2.75), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Criticize_Score

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(10) 

name_box = slide4.shapes.add_textbox(left=Inches(5),top = Inches(2.75), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Empathy_Score

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(10) 

name_box = slide4.shapes.add_textbox(left=Inches(7),top = Inches(2.75), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Search_Score 

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(10) 

#Slide 5

slide5 = prs.slides.add_slide(blank_slide_layout)
#slide1.width = Inches(8.27)
#slide1.height = Inches(11.69)

img_path5 = "Slide5img.jpg"
left = Inches(0)
top = Inches(0)

pic5 = slide5.shapes.add_picture(img_path5, left, top)

name_box = slide5.shapes.add_textbox(left=Inches(1.5),top = Inches(1.69), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Advice_score 

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(12)

name_box = slide5.shapes.add_textbox(left=Inches(1.5),top = Inches(4.35), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Criticize_Score

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(12)

name_box = slide5.shapes.add_textbox(left=Inches(1.5),top = Inches(6.975), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Empathy_Score 

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(12)

name_box = slide5.shapes.add_textbox(left=Inches(1.5),top = Inches(9.201), width=Inches(1.5), height=Inches(0.6))
name_frame = name_box.text_frame
name_frame.clear()
p = name_frame.paragraphs[0]
run = p.add_run()
run.text = Search_Score 

font = run.font
font.name = "Arial"
font.bold = True
font.italic = False
font.size = Pt(12)

def advice(dataa): 
    if dataa <= 4:
        scorea = sh2['D2'].value 
        #print(scorea)
        left1 = Inches(0.055)
        top1 = Inches(2.7)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorea)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        
        img_path_dot = "dot.png"
        left = Inches(4.77)
        top = Inches(2.53)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        #name_frame.alignment = 4
        
        #name_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        #name_frame.alignment = PP_ALIGN.LEFT
        #font = run.font
        
        #font.bold = True
        #font.italic = False
        #font.size = Pt(10) 
        #name_frame.text_frame.word_wrap = True

        #name_frame.clear()
        #p = name_frame.paragraphs[0]
        #run = p.add_run()
        #run.text = scorea 

        #font = run.font
        #font.name = "Arial"
        #font.bold = True
        #font.italic = False
        #font.size = Pt(10) 
          
    elif dataa <= 10:
        scorea = sh2['D3'].value
        left1 = Inches(0.055)
        top1 = Inches(2.7)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorea)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        
        img_path_dot = "dot.png"
        left = Inches(1.79)
        top = Inches(2.53)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        #print(scorea)
    elif dataa <= 19:
        scorea = sh2['D4'].value
        #print(scorea)
        left1 = Inches(0.055)
        top1 = Inches(2.7)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorea)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(3.29)
        top = Inches(2.53)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    else: 
        scorea = sh2['D5'].value
        #print(scorea)
        left1 = Inches(0.055)
        top1 = Inches(2.7)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorea)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(6.2)
        top = Inches(2.53)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
def criticize(datac): 
    if datac <= 3:
        scorec = sh2['F2'].value
        #print(scorec)
        left1 = Inches(0.055)
        top1 = Inches(5.47)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorec)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(1.79)
        top = Inches(5.25)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    elif datac <= 10:
        scorec = sh2['F3'].value
        #print(scorec)
        left1 = Inches(0.055)
        top1 = Inches(5.47)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorec)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(3.29)
        top = Inches(5.25)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    elif datac <= 19:
        scorec = sh2['F4'].value
        #print(scorec)
        left1 = Inches(0.055)
        top1 = Inches(5.47)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorec)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(4.77)
        top = Inches(5.25)
        
        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    else: 
        scorec = sh2['F5'].value
        #print(scorec)
        left1 = Inches(0.055)
        top1 = Inches(5.47)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scorec)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        
        img_path_dot = "dot.png"
        left = Inches(6.2)
        top = Inches(5.25)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
                
        
def empathy(datae): 
    if datae <= 8:
        scoree = sh2['H2'].value
        #print(scoree)
        left1 = Inches(0.055)
        top1 = Inches(7.97)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scoree)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(6.2)
        top = Inches(7.79)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    elif datae <= 14:
        scoree = sh2['H3'].value
        #print(scoree)
        left1 = Inches(0.055)
        top1 = Inches(7.97)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scoree)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(3.29)
        top = Inches(7.79)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    elif datae <= 20:
        scoree = sh2['H4'].value
        #print(scoree)
        left1 = Inches(0.055)
        top1 = Inches(7.97)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scoree)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        
        img_path_dot = "dot.png"
        left = Inches(1.79)
        top = Inches(7.79)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    else: 
        scoree = sh2['H5'].value
        #print(scoree)
        left1 = Inches(0.055)
        top1 = Inches(7.97)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scoree)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        
        img_path_dot = "dot.png"
        left = Inches(4.77)
        top = Inches(7.79)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
        
                
def search(datas): 
    if datas <= 8:
        scores = sh2['J2'].value
        #print(scores)
        left1 = Inches(0.055)
        top1 = Inches(10.2)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scores)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(6.2)
        top = Inches(9.99)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    elif datas <= 14:
        scores = sh2['J3'].value
        #print(scores)
        left1 = Inches(0.055)
        top1 = Inches(10.2)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scores)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY 
        
        img_path_dot = "dot.png"
        left = Inches(3.29)
        top = Inches(9.99)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    elif datas <= 20:
        scores = sh2['J4'].value
        #print(scores)
        left1 = Inches(0.055)
        top1 = Inches(10.2)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scores)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        
        img_path_dot = "dot.png"
        left = Inches(1.79)
        top = Inches(9.99)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
    else: 
        scores = sh2['J5'].value
        #print(scores)
        left1 = Inches(0.055)
        top1 = Inches(10.2)
        width1 = Inches(8.1)
        height1 = Inches(1)
        textbox = slide5.shapes.add_textbox(left1, top1, width1, height1)
        long_text = str(scores)
        paragraph = textbox.text_frame.add_paragraph()
        #paragraph = name_frame.add_paragraph()
        
        paragraph.text = long_text
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        
        img_path_dot = "dot.png"
        left = Inches(4.77)
        top = Inches(9.99)

        picdot = slide5.shapes.add_picture(img_path_dot, left, top)
        
        
       
advice(dataa)
criticize(datac)
empathy(datae)
search(datas)

slide6 = prs.slides.add_slide(blank_slide_layout)
#slide1.width = Inches(8.27)
#slide1.height = Inches(11.69)

img_path6 = "Slide6img.jpg"
left = Inches(0)
top = Inches(0)

pic6 = slide6.shapes.add_picture(img_path6, left, top)

# Save the changes to the PowerPoint file
prs.save(Name + " CPA.pptx")
